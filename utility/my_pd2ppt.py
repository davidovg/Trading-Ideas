import six
import pandas as pd
from math import *
from pptx import Presentation
from pptx.shapes.placeholder import TablePlaceholder
from pptx.oxml import CT_GraphicalObjectFrame
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR_INDEX
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
import io
from PIL import ImageFont
from pptx.table import Table
from pptx.text.fonts import FontFiles
import numpy as np
from pptx.oxml.shapes.graphfrm import CT_GraphicalObjectFrame
from pptx.spec import GRAPHIC_DATA_URI_TABLE
from pptx.oxml.table import CT_Table


table_styles = pd.DataFrame(data=
                      [{"style_id": "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}", "styleName": "Medium Style 2 - Accent 1"},
                      {"style_id": "{2D5ABB26-0587-4C30-8999-92F81FD0307C}", "styleName": "No Style, No Grid"},
                      {"style_id": "{3B4B98B0-60AC-42C2-AFA5-B58CD77FA1E5}", "styleName": "Light Style 1 - Accent 1"},
                      {"style_id": "{0E3FDE45-AF77-4B5C-9715-49D594BDF05E}", "styleName": "Light Style 1 - Accent 2"},
                      {"style_id": "{68D230F3-CF80-4859-8CE7-A43EE81993B5}", "styleName": "Light Style 1 - Accent 6"},
                      {"style_id": "{9D7B26C5-4107-4FEC-AEDC-1716B250A1EF}", "styleName": "Light Style 1"},
                      {"style_id": "{69012ECD-51FC-41F1-AA8D-1B2483CD663E}", "styleName": "Light Style 2 - Accent 1"},
                      {"style_id": "{8A107856-5554-42FB-B03E-39F5DBC370BA}", "styleName": "Medium Style 4 - Accent 2"},
                      {"style_id": "{08FB837D-C827-4EFA-A057-4D05807E0F7C}", "styleName": "Themed Style 1 - Accent 6"},
                      {"style_id": "{21E4AEA4-8DFA-4A89-87EB-49C32662AFE0}", "styleName": "Medium Style 2 - Accent 2"},
                      {"style_id": "{93296810-A885-4BE3-A3E7-6D5BEEA58F35}", "styleName": "Medium Style 2 - Accent 6"},
                      {"style_id": "{5A111915-BE36-4E01-A7E5-04B1672EAD32}", "styleName": "Light Style 2 - Accent 5"},
                      {"style_id": "{9DCAF9ED-07DC-4A11-8D7F-57B35C25682E}", "styleName": "Medium Style 1 - Accent 2"},
                      {"style_id": "{B301B821-A1FF-4177-AEE7-76D212191A09}", "styleName": "Medium Style 1 - Accent 1"},
                      {"style_id": "{1E171933-4619-4E11-9A3F-F7608DF75F80}", "styleName": "Medium Style 1 - Accent 4"},
                      {"style_id": "{10A1B5D5-9B99-4C35-A422-299274C87663}", "styleName": "Medium Style 1 - Accent 6"},
                      {"style_id": "{793D81CF-94F2-401A-BA57-92F5A7B2D0C5}", "styleName": "Medium Style 1"}
                       ])

style_0 = table_styles["style_id"].loc[0]
style_1 = table_styles["style_id"].loc[1]

default_margin = [0,0,0,0]
round_to_n = lambda x, n: round(x, -int(floor(log10(abs(x)))) + (n - 1))

_color_dict = {"AliceBlue": RGBColor(0xF0, 0xF8, 0xFF), "AntiqueWhite": RGBColor(0xFA, 0xEB, 0xD7),
     "Aqua": RGBColor(0x00, 0xFF, 0xFF), "Aquamarine": RGBColor(0x7F, 0xFF, 0xD4), "Azure": RGBColor(0xF0, 0xFF, 0xFF),
     "Beige": RGBColor(0xF5, 0xF5, 0xDC), "Bisque": RGBColor(0xFF, 0xE4, 0xC4), "Black": RGBColor(0x00, 0x00, 0x00),
     "BlanchedAlmond": RGBColor(0xFF, 0xEB, 0xCD), "Blue": RGBColor(0x00, 0x00, 0xFF),
     "BlueViolet": RGBColor(0x8A, 0x2B, 0xE2), "Brown": RGBColor(0xA5, 0x2A, 0x2A),
     "BurlyWood": RGBColor(0xDE, 0xB8, 0x87), "CadetBlue": RGBColor(0x5F, 0x9E, 0xA0),
     "Chartreuse": RGBColor(0x7F, 0xFF, 0x00), "Chocolate": RGBColor(0xD2, 0x69, 0x1E),
     "Coral": RGBColor(0xFF, 0x7F, 0x50), "CornflowerBlue": RGBColor(0x64, 0x95, 0xED),
     "Cornsilk": RGBColor(0xFF, 0xF8, 0xDC), "Crimson": RGBColor(0xDC, 0x14, 0x3C), "Cyan": RGBColor(0x00, 0xFF, 0xFF),
     "DarkBlue": RGBColor(0x00, 0x00, 0x8B), "DarkCyan": RGBColor(0x00, 0x8B, 0x8B),
     "DarkGoldenRod": RGBColor(0xB8, 0x86, 0x0B), "DarkGray": RGBColor(0xA9, 0xA9, 0xA9),
     "DarkGreen": RGBColor(0x00, 0x64, 0x00), "DarkKhaki": RGBColor(0xBD, 0xB7, 0x6B),
     "DarkMagenta": RGBColor(0x8B, 0x00, 0x8B), "DarkOliveGreen": RGBColor(0x55, 0x6B, 0x2F),
     "DarkOrange": RGBColor(0xFF, 0x8C, 0x00), "DarkOrchid": RGBColor(0x99, 0x32, 0xCC),
     "DarkRed": RGBColor(0x8B, 0x00, 0x00), "DarkSalmon": RGBColor(0xE9, 0x96, 0x7A),
     "DarkSeaGreen": RGBColor(0x8F, 0xBC, 0x8F), "DarkSlateBlue": RGBColor(0x48, 0x3D, 0x8B),
     "DarkSlateGray": RGBColor(0x2F, 0x4F, 0x4F), "DarkTurquoise": RGBColor(0x00, 0xCE, 0xD1),
     "DarkViolet": RGBColor(0x94, 0x00, 0xD3), "DeepPink": RGBColor(0xFF, 0x14, 0x93),
     "DeepSkyBlue": RGBColor(0x00, 0xBF, 0xFF), "DimGray": RGBColor(0x69, 0x69, 0x69),
     "DodgerBlue": RGBColor(0x1E, 0x90, 0xFF), "FireBrick": RGBColor(0xB2, 0x22, 0x22),
     "FloralWhite": RGBColor(0xFF, 0xFA, 0xF0), "ForestGreen": RGBColor(0x22, 0x8B, 0x22),
     "Fuchsia": RGBColor(0xFF, 0x00, 0xFF), "Gainsboro": RGBColor(0xDC, 0xDC, 0xDC),
     "GhostWhite": RGBColor(0xF8, 0xF8, 0xFF), "Gold": RGBColor(0xFF, 0xD7, 0x00),
     "GoldenRod": RGBColor(0xDA, 0xA5, 0x20), "Gray": RGBColor(0x80, 0x80, 0x80), "Green": RGBColor(0x00, 0x80, 0x00),
     "GreenYellow": RGBColor(0xAD, 0xFF, 0x2F), "HoneyDew": RGBColor(0xF0, 0xFF, 0xF0),
     "HotPink": RGBColor(0xFF, 0x69, 0xB4), "IndianRed": RGBColor(0xCD, 0x5C, 0x5C),
     "Indigo": RGBColor(0x4B, 0x00, 0x82), "Ivory": RGBColor(0xFF, 0xFF, 0xF0), "Khaki": RGBColor(0xF0, 0xE6, 0x8C),
     "Lavender": RGBColor(0xE6, 0xE6, 0xFA), "LavenderBlush": RGBColor(0xFF, 0xF0, 0xF5),
     "LawnGreen": RGBColor(0x7C, 0xFC, 0x00), "LemonChiffon": RGBColor(0xFF, 0xFA, 0xCD),
     "LightBlue": RGBColor(0xAD, 0xD8, 0xE6), "LightCoral": RGBColor(0xF0, 0x80, 0x80),
     "LightCyan": RGBColor(0xE0, 0xFF, 0xFF), "LightGoldenRodYellow": RGBColor(0xFA, 0xFA, 0xD2),
     "LightGray": RGBColor(0xD3, 0xD3, 0xD3), "LightGreen": RGBColor(0x90, 0xEE, 0x90),
     "LightPink": RGBColor(0xFF, 0xB6, 0xC1), "LightSalmon": RGBColor(0xFF, 0xA0, 0x7A),
     "LightSeaGreen": RGBColor(0x20, 0xB2, 0xAA), "LightSkyBlue": RGBColor(0x87, 0xCE, 0xFA),
     "LightSlateGray": RGBColor(0x77, 0x88, 0x99), "LightSteelBlue": RGBColor(0xB0, 0xC4, 0xDE),
     "LightYellow": RGBColor(0xFF, 0xFF, 0xE0), "Lime": RGBColor(0x00, 0xFF, 0x00),
     "LimeGreen": RGBColor(0x32, 0xCD, 0x32), "Linen": RGBColor(0xFA, 0xF0, 0xE6),
     "Magenta": RGBColor(0xFF, 0x00, 0xFF), "Maroon": RGBColor(0x80, 0x00, 0x00),
     "MediumAquaMarine": RGBColor(0x66, 0xCD, 0xAA), "MediumBlue": RGBColor(0x00, 0x00, 0xCD),
     "MediumOrchid": RGBColor(0xBA, 0x55, 0xD3), "MediumPurple": RGBColor(0x93, 0x70, 0xDB),
     "MediumSeaGreen": RGBColor(0x3C, 0xB3, 0x71), "MediumSlateBlue": RGBColor(0x7B, 0x68, 0xEE),
     "MediumSpringGreen": RGBColor(0x00, 0xFA, 0x9A), "MediumTurquoise": RGBColor(0x48, 0xD1, 0xCC),
     "MediumVioletRed": RGBColor(0xC7, 0x15, 0x85), "MidnightBlue": RGBColor(0x19, 0x19, 0x70),
     "MintCream": RGBColor(0xF5, 0xFF, 0xFA), "MistyRose": RGBColor(0xFF, 0xE4, 0xE1),
     "Moccasin": RGBColor(0xFF, 0xE4, 0xB5), "NavajoWhite": RGBColor(0xFF, 0xDE, 0xAD),
     "Navy": RGBColor(0x00, 0x00, 0x80), "OldLace": RGBColor(0xFD, 0xF5, 0xE6), "Olive": RGBColor(0x80, 0x80, 0x00),
     "OliveDrab": RGBColor(0x6B, 0x8E, 0x23), "Orange": RGBColor(0xFF, 0xA5, 0x00),
     "OrangeRed": RGBColor(0xFF, 0x45, 0x00), "Orchid": RGBColor(0xDA, 0x70, 0xD6),
     "PaleGoldenRod": RGBColor(0xEE, 0xE8, 0xAA), "PaleGreen": RGBColor(0x98, 0xFB, 0x98),
     "PaleTurquoise": RGBColor(0xAF, 0xEE, 0xEE), "PaleVioletRed": RGBColor(0xDB, 0x70, 0x93),
     "PapayaWhip": RGBColor(0xFF, 0xEF, 0xD5), "PeachPuff": RGBColor(0xFF, 0xDA, 0xB9),
     "Peru": RGBColor(0xCD, 0x85, 0x3F), "Pink": RGBColor(0xFF, 0xC0, 0xCB), "Plum": RGBColor(0xDD, 0xA0, 0xDD),
     "PowderBlue": RGBColor(0xB0, 0xE0, 0xE6), "Purple": RGBColor(0x80, 0x00, 0x80), "Red": RGBColor(0xFF, 0x00, 0x00),
     "RosyBrown": RGBColor(0xBC, 0x8F, 0x8F), "RoyalBlue": RGBColor(0x41, 0x69, 0xE1),
     "SaddleBrown": RGBColor(0x8B, 0x45, 0x13), "Salmon": RGBColor(0xFA, 0x80, 0x72),
     "SandyBrown": RGBColor(0xF4, 0xA4, 0x60), "SeaGreen": RGBColor(0x2E, 0x8B, 0x57),
     "SeaShell": RGBColor(0xFF, 0xF5, 0xEE), "Sienna": RGBColor(0xA0, 0x52, 0x2D), "Silver": RGBColor(0xC0, 0xC0, 0xC0),
     "SkyBlue": RGBColor(0x87, 0xCE, 0xEB), "SlateBlue": RGBColor(0x6A, 0x5A, 0xCD),
     "SlateGray": RGBColor(0x70, 0x80, 0x90), "Snow": RGBColor(0xFF, 0xFA, 0xFA),
     "SpringGreen": RGBColor(0x00, 0xFF, 0x7F), "SteelBlue": RGBColor(0x46, 0x82, 0xB4),
     "Tan": RGBColor(0xD2, 0xB4, 0x8C), "Teal": RGBColor(0x00, 0x80, 0x80), "Thistle": RGBColor(0xD8, 0xBF, 0xD8),
     "Tomato": RGBColor(0xFF, 0x63, 0x47), "Turquoise": RGBColor(0x40, 0xE0, 0xD0),
     "Violet": RGBColor(0xEE, 0x82, 0xEE), "Wheat": RGBColor(0xF5, 0xDE, 0xB3), "White": RGBColor(0xFF, 0xFF, 0xFF),
     "WhiteSmoke": RGBColor(0xF5, 0xF5, 0xF5), "Yellow": RGBColor(0xFF, 0xFF, 0x00),
     "YellowGreen": RGBColor(0x9A, 0xCD, 0x32), "EquityGreen": RGBColor(0x92, 0xD0, 0x50),
     "BondBlue": RGBColor(0x00, 0x20, 0x60), "AltGold": RGBColor(0xFF, 0xC0,0x00),
     "CashGray": RGBColor(0xA6, 0xA6, 0xA6)}


def _do_formatting(value, format_str):
    """Format value according to format_str, and deal
    sensibly with format_str if it is missing or invalid.
    """
    if format_str == '':
        if type(value) in six.integer_types:
            format_str = ','
        elif type(value) is float:
            format_str = 'f'
        elif type(value) is str:
            format_str = 's'
    elif format_str[0] == '.':
        if format_str.endswith('R'):
            if type(value) in six.integer_types:
                value = round_to_n(value, int(format_str[1]))
                format_str = ','
        if not format_str.endswith('G'):
            format_str = format_str + "G"
    try:
        value = format(value, format_str)
    except:
        value = format(value, '')

    return value


def process_position_parameter(param):
    """Process positioning parameters (left, top, width, height) given to
    df_to_table.

    If an integer, returns the right instance of the Cm class to allow it to be
    treated as cm. If missing, then default to 4cm. Otherwise, pass through
    whatever it gets.
    """
    if param is None:
        return Cm(4)
    elif type(param) is int:
        return Cm(param)
    else:
        return param


def df_to_table(slide, df, left=None, top=None, width=None, height=None,
                colnames=None, col_formatters=None, rounding=None,
                name=None):
    """Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation.

    The table is a standard Powerpoint table, and can easily be modified with
    the Powerpoint tools, for example: resizing columns, changing formatting etc.

    Parameters
    ----------
    slide: ``pptx.slide.Slide``
        slide object from the python-pptx library containing the slide on which
        you want the table to appear

    df: pandas ``DataFrame``
       DataFrame with the data

    left: int, optional
       Position of the left-side of the table, either as an integer in cm, or
       as an instance of a pptx.util Length class (pptx.util.Inches for
       example). Defaults to 4cm.

    top: int, optional
       Position of the top of the table, takes parameters as above.

    width: int, optional
       Width of the table, takes parameters as above.

    height: int, optional
       Height of the table, takes parameters as above.

    col_formatters: list, optional
       A n_columns element long list containing format specifications for each
       column. For example ['', ',', '.2'] does no special formatting for the
       first column, uses commas as thousands separators in the second column,
       and formats the third column as a float with 2 decimal places.

    rounding: list, optional
       A n_columns element long list containing a number for each integer
       column that requires rounding that is then multiplied by -1 and passed
       to round(). The practical upshot of this is that you can give something
       like ['', 3, ''], which does nothing for the 1st and 3rd columns (as
       they aren't integer values), but for the 2nd column, rounds away the 3
       right-hand digits (eg. taking 25437 to 25000).

    name: str, optional
       A name to be given to the table in the Powerpoint file. This is not
       displayed, but can help extract the table later to make further changes.

    Returns
    -------
    pptx.shapes.graphfrm.GraphicFrame
        The python-pptx table (GraphicFrame) object that was created (which can
        then be used to do further manipulation if desired)
    """
    left = process_position_parameter(left)
    top = process_position_parameter(top)
    width = process_position_parameter(width)
    height = process_position_parameter(height)

    rows, cols = df.shape
    shp = slide.shapes.add_table(rows+1, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        shp.table.cell(0, col_index).text = col_name

    m = df.values

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]

            if col_formatters is None:
                text = str(val)
            else:
                text = _do_formatting(val, col_formatters[col])

            shp.table.cell(row+1, col).text = text

    if name is not None:
        shp.name = name

    return shp


def df_to_table_placeholder(placeholder, df, colnames=None, col_formatters=None,
                name=None):
    """My variation of the function df_to_table that insert the table in a placeholder
    """

    rows, cols = df.shape
    shp = placeholder.insert_table(rows+1, cols)

    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        shp.table.cell(0, col_index).text = col_name

    m = df.values

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]

            if col_formatters is None:
                text = str(val)
            else:
                text = _do_formatting(val, col_formatters[col])

            shp.table.cell(row+1, col).text = text

    if name is not None:
        shp.name = name

    return shp


def auto_uniform_format_table(table: Table, font_size=10, font_name='Cambria', margins=default_margin):
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            cell.text_frame.word_wrap = False
            cell.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            cell.margin_left = 0
            cell.margin_top = 0
            cell.margin_bottom = 0
            cell.margin_right = 0

            cell.text_frame.margin_left = margins[0]
            cell.text_frame.margin_top = margins[1]
            cell.text_frame.margin_bottom = margins[2]
            cell.text_frame.margin_right = margins[3]

            # cell.text_frame.margin_bottom = 0
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.name = font_name
                    run.font.color.theme_color = MSO_THEME_COLOR_INDEX.DARK_2
                    if row_idx == 0:
                        run.font.bold = True
                        run.font.italic = True
                        run.font.color.theme_color = MSO_THEME_COLOR_INDEX.LIGHT_1


def autoset_column_width(table, df, pts, font="Cambria"):
    j = 0

    # TODO: It might happen that the is_bold, is_italics combination you arre searching does not exist
    # TODO: INSERT error handling for both missing family_name (real error) and missing b/i
    font = ImageFont.truetype(FontFiles.find(font, True, False), size=pts)
    # font = ImageFont.truetype("C:Windows/Fonts/cambriaz.ttf", size=pts)
    width_font = lambda x: font.getsize(str(x))[0]
    old_width = 0.0
    new_width = 0.0
    for i in range(len(df.columns)):
        m = max(df.iloc[:,i].apply(width_font).max(), width_font(df.columns[i]))
        old_width = old_width + table.columns[j].width
        table.columns[j].width = Pt(m+5)
        new_width = new_width + table.columns[j].width
        j = j+1

    mult = old_width/new_width
    #Final rescaling
    for i in range(len(df.columns)-1):
        val = int(table.columns[i].width * mult)
        old_width = old_width - val
        table.columns[i].width = val

    table.columns[len(df.columns)-1].width = int(old_width)


def autoset_row_height(table, df, pts, font="Cambria", nhf=False, max_rescale_height=None):
    j = 0

    # TODO: It might happen that the is_bold, is_italics combination you arre searching does not exist
    # TODO: INSERT error handling for both missing family_name (real error) and missing b/i
    font = ImageFont.truetype(FontFiles.find(font, True, False), size=pts)
    # font = ImageFont.truetype("C:Windows/Fonts/cambriaz.ttf", size=pts)
    height_font = np.vectorize(lambda x: font.getsize(str(x))[1])
    old_height = 0.0
    new_height = 0.0

    m = max(height_font(df.columns.values))
    old_height = old_height + table.rows[j].height
    table.rows[j].height = Pt(m + 2)
    new_height = new_height + table.rows[j].height
    j = j + 1

    for _k, _r in df.iterrows():
        m = max(height_font(_r.values))
        old_height = old_height + table.rows[j].height
        table.rows[j].height = Pt(m+2)
        new_height = new_height + table.rows[j].height
        j = j+1

    if not nhf:
        mult = old_height/new_height if max_rescale_height is None \
                                    else min(old_height/new_height, max_rescale_height)
        old_height = mult * new_height
        #Final rescaling
        for i in range(1,len(df)+1):
            val = int(table.rows[i].height * mult)
            old_height = old_height - val
            table.rows[i].height = val

        table.rows[0].height = int(old_height)


def auto_df_to_table_plc(placeholder, df, colnames=None, col_formatters=None, name=None,
                         font_size=8, style=style_0, no_height_fit=False, max_rescale_height=None):
    if isinstance(placeholder, TablePlaceholder):
        rows, cols = df.shape
        placeholder = MyTablePlaceholder(placeholder, style)

        shp = placeholder.insert_table(rows + 1, cols)

        if colnames is None:
            colnames = list(df.columns)

        # Insert the column names
        for col_index, col_name in enumerate(colnames):
            shp.table.cell(0, col_index).text = col_name

        m = df.values

        for row in range(rows):
            for col in range(cols):
                val = m[row, col]

                if col_formatters is None:
                    text = str(val)
                else:
                    text = _do_formatting(val, col_formatters[col])

                shp.table.cell(row + 1, col).text = text

        if name is not None:
            shp.name = name

        table = shp.table

        auto_uniform_format_table(table, font_size, "Cambria")
        autoset_column_width(table, df, font_size, "Cambria")
        autoset_row_height(table, df, font_size, "Cambria", no_height_fit, max_rescale_height)

        return shp


def df_to_powerpoint(filename, df, **kwargs):
    """Converts a Pandas DataFrame to a table in a new, blank PowerPoint
    presentation.

    Creates a new PowerPoint presentation with the given filename, with a single
    slide containing a single table with the Pandas DataFrame data in it.

    The table is a standard Powerpoint table, and can easily be modified with
    the Powerpoint tools, for example: resizing columns, changing formatting
    etc.

    Parameters
    ----------
    filename: Filename to save the PowerPoint presentation as

    df: pandas ``DataFrame``
        DataFrame with the data

    **kwargs
        All other arguments that can be taken by ``df_to_table()`` (such as
        ``col_formatters`` or ``rounding``) can also be passed here.

    Returns
    -------
    pptx.shapes.graphfrm.GraphicFrame
        The python-pptx table (GraphicFrame) object that was created (which can
        then be used to do further manipulation if desired)
    """
    pres = Presentation()
    blank_slide_layout = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank_slide_layout)
    table = df_to_table(slide, df, **kwargs)
    pres.save(filename)

    return table


def uniform_format_table(table, font_size=10):
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.name = 'Cambria'
                    run.font.color.theme_color = MSO_THEME_COLOR_INDEX.DARK_2
                    if row_idx == 0:
                        run.font.bold = True
                        run.font.italic = True
                        run.font.color.theme_color = MSO_THEME_COLOR_INDEX.LIGHT_1


def from_fig_to_fileLike(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png')
    buf.seek(0)
    return buf


def get_table_width(table):
    output = 0
    for column in table.columns:
        output = output + column.width
    return output


def get_row_height(table, row_id):
    output = table.cell(row_id, 1).height
    return output


def set_row_height_floor(table, flr):
    for row_idx, row in enumerate(table.rows):
        row.height = min(row.height, flr)


def set_first_col_width(table, dim):
    table_width = get_table_width(table)
    for i in len(table.columns):
        if i == 0:
            table.columns[0].width = Cm(dim)
        else:
            table.columns[i].width = Cm((table_width - dim)/(len(table.columns)-1))


def set_column_width(table, dic):
    for i in dic:
        table.columns[int(i)].width = Cm(dic[i])


def color_single_row(table, row_numb=0, color=RGBColor(0xFB, 0x8F, 0x00)):
    row = table.rows[row_numb]
    for col_idx, cell in enumerate(row.cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = color


def color_single_column(table, col_numb=0, color=RGBColor(203, 213, 232)):
    for row_id in range(1, len(table.rows)):
        table.cell(row_id, col_numb).fill.solid()
        table.cell(row_id, col_numb).fill.fore_color.rgb = color



class CT_GraphicalObjectFrame_TableStyle(CT_GraphicalObjectFrame):

    @classmethod
    def new_table_graphicFrame_TableStyle(cls, id_, name, rows, cols, x, y, cx, cy, ts):
        graphicFrame = cls.new_graphicFrame(id_, name, x, y, cx, cy)
        graphicFrame.graphic.graphicData.uri = GRAPHIC_DATA_URI_TABLE
        graphicFrame.graphic.graphicData.append(CT_Table.new_tbl(rows, cols, cx, cy, ts))
        return graphicFrame

class MyTablePlaceholder(TablePlaceholder):

    def __init__(self, old_t, _ts=style_0):
        self.__class__ = type(old_t.__class__.__name__,
                              (self.__class__, old_t.__class__),
                              {})
        self.__dict__ = old_t.__dict__
        self.ts = _ts

    def _new_placeholder_table(self, rows, cols):
        """
        Return a newly added `p:graphicFrame` element containing an empty
        table with *rows* rows and *cols* columns, positioned at the location
        of this placeholder and having its same width. The table's height is
        determined by the number of rows.
        """

        # TODO: Set style in pptx.oxml.table line 102

        shape_id, name = self.shape_id, self.name
        return CT_GraphicalObjectFrame_TableStyle.new_table_graphicFrame_TableStyle(
            shape_id, name, rows, cols, self.left, self.top, self.width, self.height, self.ts
        )
        # shape_id, name = self.shape_id, self.name
        # return CT_GraphicalObjectFrame.new_table_graphicFrame(
        #     shape_id, name, rows, cols, self.left, self.top, self.width, self.height
        # )


def my_colors(color_name):
    return _color_dict.get(color_name, _color_dict["DarkRed"])


def resize_mplfig_on_plc(fig, plc):
    h = plc.height.inches
    w = plc.width.inches
    fig.set_size_inches(w, h)

def value_cond_format(table, lower, upper, position="outside", colour="DarkOrange"):
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            if row_idx > 0 and col_idx > 0:
                if position == "outside":
                    if (float(cell.text_frame.text) < lower or float(cell.text_frame.text) > upper):
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = my_colors(colour)
                else:
                    if (float(cell.text_frame.text) > lower and float(cell.text_frame.text) < upper):
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = my_colors(colour)

def merge_cell(table, top_left_corner, bottom_right_corner):
    n_rows = range(top_left_corner[0], bottom_right_corner[0]+1)
    n_cols = range(top_left_corner[1], bottom_right_corner[1]+1)
    for i in n_rows:
        for j in n_cols:
            k = (i != top_left_corner[0] or j != top_left_corner[1])
            if k:
                table.cell(i, j).text_frame.clear()
    table.cell(top_left_corner[0], top_left_corner[1]).merge(table.cell(bottom_right_corner[0], bottom_right_corner[1]))